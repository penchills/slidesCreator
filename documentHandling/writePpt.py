from pptx import Presentation

class pptHandling:
    def __init__(self, file_path):
        self.file_path = file_path
        self.presentation = None
    
    def create_presentation(self):
        """Create a new PowerPoint presentation."""
        self.presentation = Presentation()
        return self.presentation
    
    def add_slide(self, title, content):
        """Add a slide with a title and content to the presentation."""
        if self.presentation is None:
            raise ValueError("Presentation not created. Call create_presentation() first.")
        
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.shapes.title
        content_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        content_placeholder.text = content
        return slide
    
    def save_presentation(self):
        """Save the PowerPoint presentation to the specified file path."""
        if self.presentation is None:
            raise ValueError("Presentation not created. Call create_presentation() first.")
        
        self.presentation.save(self.file_path)
        return self.file_path
    
    def load_presentation(self):
        """Load an existing PowerPoint presentation from the specified file path."""
        if not self.file_path:
            raise ValueError("File path not specified.")
        
        self.presentation = Presentation(self.file_path)
        return self.presentation
